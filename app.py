import io, re, json, uuid
from pathlib import Path
from typing import List
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup

BASE=Path(__file__).resolve().parent
REPORTS=BASE/'reports'; REPORTS.mkdir(exist_ok=True)
CRITERIA=json.loads((BASE/'institutional_criteria.json').read_text(encoding='utf-8'))
app=FastAPI(title='UCAN Course Evaluator')

def extract(name,raw):
    ext=Path(name).suffix.lower()
    if ext=='.pdf':
        r=PdfReader(io.BytesIO(raw)); return '\n'.join((p.extract_text() or '') for p in r.pages)
    if ext=='.docx':
        d=Document(io.BytesIO(raw)); out=[p.text for p in d.paragraphs if p.text.strip()]
        for t in d.tables:
            for row in t.rows: out.append(' | '.join(c.text for c in row.cells))
        return '\n'.join(out)
    if ext=='.pptx':
        prs=Presentation(io.BytesIO(raw)); out=[]
        for i,s in enumerate(prs.slides,1):
            out.append(f'[SLIDE {i}]')
            for sh in s.shapes:
                if hasattr(sh,'text') and sh.text.strip(): out.append(sh.text)
        return '\n'.join(out)
    if ext=='.xlsx':
        wb=load_workbook(io.BytesIO(raw),data_only=True,read_only=True); out=[]
        for ws in wb.worksheets:
            out.append(f'[SHEET {ws.title}]')
            for row in ws.iter_rows(values_only=True):
                vals=[str(v) for v in row if v is not None]
                if vals: out.append(' | '.join(vals))
        return '\n'.join(out)
    if ext in {'.html','.htm'}: return BeautifulSoup(raw.decode('utf-8','ignore'),'html.parser').get_text('\n')
    if ext in {'.txt','.md','.csv'}: return raw.decode('utf-8','ignore')
    return f'[FORMATO NO EXTRAÍDO AUTOMÁTICAMENTE: {name}]'

def finding(criterion,status,priority,location,evidence,comment,recommendation,example=''):
    return dict(criterion=criterion,status=status,priority=priority,location=location,evidence=evidence,comment=comment,recommendation=recommendation,example=example)

def contains(text,words): return any(w.lower() in text.lower() for w in words)
def count(text,patterns): return sum(len(re.findall(p,text,re.I|re.M)) for p in patterns)

def evaluate_rules(meta, docs):
    alltext='\n'.join(x['text'] for x in docs); low=alltext.lower(); names=', '.join(x['name'] for x in docs)
    syllabus=[]; inst=[]; qm=[]; udl=[]; five=[]; objectives=[]; citations=[]; compliance=[]; modules=[]; actions=[]
    has_syllabus=contains(low,['prontuario','syllabus']); has_plan=contains(low,['plan de trabajo','course schedule','calendar']); has_intro=contains(low,['introducción','introduction','bienvenida','welcome'])
    syllabus.append(finding('Prontuario oficial','Meets' if has_syllabus else 'Does Not Meet','High',names,'Se detectó prontuario/syllabus.' if has_syllabus else 'No se detectó un prontuario identificable.','El prontuario es base de la alineación curricular.','Incluir el prontuario oficial y actualizado.'))
    syllabus.append(finding('Plan de trabajo','Meets' if has_plan else 'Does Not Meet','High',names,'Se detectó plan/calendario.' if has_plan else 'No se detectó plan de trabajo.','Debe documentar la secuencia temporal del curso.','Añadir plan de trabajo con semanas, objetivos, contenidos, actividades y fechas.'))
    syllabus.append(finding('Introducción y orientación','Meets' if has_intro else 'Partially Meets','Medium',names,'Indicadores de introducción/bienvenida: '+str(has_intro),'La orientación inicial facilita la navegación y expectativas.','Incluir bienvenida, propósito, naturaleza del curso y pasos iniciales.'))
    # Institutional structural rules
    checks=[('Introducción al curso',['bienvenida','welcome','introducción','introduction']),('Objetivos medibles',['objetivo','objective','competencia']),('Actividades de aprendizaje',['actividad','assignment','discussion','foro']),('Evaluación y criterios',['rúbrica','rubric','criterio','assessment','evaluación']),('Interacción',['foro','discussion','feedback','retroalimentación']),('Accesibilidad',['accesibilidad','accessibility','alt text','texto alternativo','subtítulo','caption']),('Derechos de autor',['copyright','derechos de autor','creative commons','licencia'])]
    for label,terms in checks:
        ok=contains(low,terms); inst.append(finding(label,'Meets' if ok else 'Partially Meets','High' if label in ['Objetivos medibles','Evaluación y criterios','Accesibilidad'] else 'Medium',names,'Términos/evidencia detectados: '+(', '.join([t for t in terms if t in low]) or 'ninguno'),'Revisión automatizada estructural.','Verificar manualmente la calidad y alineación; añadir evidencia explícita si falta.'))
    # QM high-level
    for label,terms in [('Course Overview',['bienvenida','welcome','prontuario','syllabus']),('Learning Objectives',['objetivo','objective']),('Assessment and Measurement',['rúbrica','rubric','assessment','evaluación']),('Instructional Materials',['referencias','references','lectura','reading']),('Learning Activities and Interaction',['actividad','foro','discussion']),('Course Technology',['blackboard','lms','tecnología','technology']),('Learner Support',['apoyo','support','cai','biblioteca','library']),('Accessibility and Usability',['accesibilidad','accessibility','alt text','caption'])]:
        ok=contains(low,terms); qm.append(finding(label,'Meets' if ok else 'Partially Meets','Medium',names,'Indicadores localizados: '+(', '.join([t for t in terms if t in low]) or 'ninguno'),'Revisión QM de alto nivel, no certificación oficial.','Revisar cualitativamente el criterio y documentar evidencia.'))
    # UDL
    for label,terms,rec in [('Engagement',['elección','choice','colabor','discussion','foro','reflex'],'Ofrecer opciones de participación, relevancia y autorregulación.'),('Representation',['video','imagen','image','audio','lectura','infografía','transcript'],'Presentar conceptos en más de un formato accesible.'),('Action and Expression',['proyecto','presentación','presentation','ensayo','video','producto'],'Permitir diversas formas de demostrar el aprendizaje cuando sea apropiado.')]:
        ok=contains(low,terms); udl.append(finding(label,'Meets' if ok else 'Partially Meets','Medium',names,'Indicadores: '+(', '.join([t for t in terms if t in low]) or 'ninguno'),'Evidencia automatizada de DUA/UDL.',' '+rec))
    # 5E
    phases=[('Engage',['pregunta','caso','diagnóstico','conocimientos previos','motivación']),('Explore',['explora','investiga','indaga','simulación','analiza datos']),('Explain',['explica','contenido','lectura','concepto','teoría']),('Elaborate',['aplica','proyecto','caso','transfer','diseña']),('Evaluate',['evalúa','rúbrica','prueba','reflexión','assessment'])]
    for phase,terms in phases:
        hits=[t for t in terms if t in low]; five.append(finding('5E – '+phase,'Meets' if hits else 'Partially Meets','Medium',names,'Indicadores: '+(', '.join(hits) if hits else 'ninguno'),'La detección confirma indicadores, pero la secuencia pedagógica debe validarse por un revisor.','Asegurar que la fase cumpla su función dentro de la secuencia Engage → Explore → Explain → Elaborate → Evaluate.'))
    # Objectives Bloom/Webb heuristic
    bloom={'Remember':['define','identify','list','recall','identifica','enumera','define'],'Understand':['explain','describe','summarize','explica','describe','resume'],'Apply':['apply','use','demonstrate','aplica','utiliza','demuestra'],'Analyze':['analyze','compare','differentiate','analiza','compara','diferencia'],'Evaluate':['evaluate','justify','critique','evalúa','justifica','critica'],'Create':['create','design','develop','crea','diseña','desarrolla']}
    objective_lines=[ln.strip() for ln in alltext.splitlines() if len(ln.strip())<350 and contains(ln,['objetivo','objective'])]
    sample=' | '.join(objective_lines[:12]) if objective_lines else 'No se pudieron aislar objetivos de forma confiable.'
    levels=[]
    for lvl,verbs in bloom.items():
        if contains(low,verbs): levels.append(lvl)
    objectives.append(finding('Taxonomía revisada de Bloom','Partially Meets','High',names,sample,'Niveles detectados por verbos: '+(', '.join(levels) if levels else 'no determinados')+'. Esta clasificación es heurística y requiere revisar el objetivo completo.','Usar verbos observables y verificar situación/condición y criterio de adecuacidad.','Ejemplo: Al finalizar el módulo, el estudiante analizará un caso utilizando los criterios establecidos en la rúbrica.'))
    dok='DOK 4' if contains(low,['investigación','research project','proyecto final']) else ('DOK 3' if contains(low,['justifica','critica','análisis de caso','case study']) else ('DOK 2' if contains(low,['compara','aplica','organiza']) else 'DOK 1'))
    objectives.append(finding('Norman Webb Depth of Knowledge','Partially Meets','High',names,'Indicador global estimado: '+dok,'DOK no debe determinarse solo por el verbo. La estimación usa señales de complejidad de las tareas detectadas.','Validar cada objetivo y assessment individualmente contra la demanda cognitiva real.'))
    # Citation checks
    intext=count(alltext,[r'\([A-ZÁÉÍÓÚÑ][A-Za-zÁÉÍÓÚÑáéíóúñ\-]+(?: et al\.)?,? \d{4}[a-z]?\)',r'[A-ZÁÉÍÓÚÑ][A-Za-zÁÉÍÓÚÑáéíóúñ\-]+ \(\d{4}[a-z]?\)'])
    refs=contains(low,['referencias','references','works cited','bibliografía'])
    citations.append(finding(meta['citation_style']+' – citas en el texto','Meets' if intext>=3 else 'Does Not Meet','High',names,f'Se detectaron aproximadamente {intext} patrones de cita autor-fecha.','El conteo es una señal estructural; no garantiza exactitud del estilo.','Revisar cada párrafo académico derivado de fuentes y añadir la cita correspondiente donde falte.'))
    citations.append(finding(meta['citation_style']+' – lista de referencias','Meets' if refs else 'Does Not Meet','High',names,'Se detectó sección de referencias.' if refs else 'No se detectó una sección de referencias.','Debe existir correspondencia entre las citas del contenido y la lista final.','Corregir y uniformar todas las referencias; verificar autor, fecha, título, fuente y DOI/URL según corresponda.'))
    # Compliance
    for label,terms,rec in [('Accesibilidad/WCAG',['alt text','texto alternativo','caption','subtítulo','transcript','encabezado'],'Validar además con WAVE/Ally y revisión manual.'),('Propiedad intelectual',['copyright','derechos de autor','creative commons','licencia'],'Verificar atribución, permiso/licencia y originalidad de materiales de terceros.')]:
        hits=[t for t in terms if t in low]; compliance.append(finding(label,'Partially Meets' if hits else 'Does Not Meet','High',names,'Indicadores: '+(', '.join(hits) if hits else 'ninguno'),'No puede certificarse técnicamente solo mediante extracción de texto.',rec))
    # Per-file/module snapshot
    for d in docs:
        text=d['text']; intro=contains(text,['introducción','introduction']); obj=contains(text,['objetivo','objective']); ref=contains(text,['referencias','references','works cited'])
        status='Meets' if intro and obj and ref else 'Partially Meets'
        modules.append(finding('Revisión estructural de '+d['name'],status,'Medium',d['name'],f'Introducción={intro}; objetivos={obj}; referencias={ref}.','Revisión automatizada del archivo.','Completar los elementos faltantes y revisar alineación, calidad académica y accesibilidad.'))
    # action plan from non-meets/partials
    combined=syllabus+inst+qm+udl+five+objectives+citations+compliance+modules
    for x in combined:
        if x['status']!='Meets': actions.append({'priority':x['priority'],'area':x['criterion'],'action':x['recommendation'],'location':x['location']})
    meets=sum(x['status']=='Meets' for x in combined); partial=sum(x['status']=='Partially Meets' for x in combined); no=sum(x['status']=='Does Not Meet' for x in combined); total=max(1,len(combined)); score=round((meets+0.5*partial)/total*100)
    recommendation='REQUIERE CORRECCIONES ANTES DE CERTIFICACIÓN' if no or partial else 'CUMPLE ESTRUCTURALMENTE; PENDIENTE VALIDACIÓN HUMANA'
    return {'summary':{'overall_score':score,'meets':meets,'partially_meets':partial,'does_not_meet':no,'executive_summary':'Evaluación automatizada basada en reglas locales y evidencia extraída de los archivos. No utiliza servicios externos de IA. Los hallazgos sirven para apoyar la revisión curricular y deben validarse por un evaluador institucional.','recommendation':recommendation},'syllabus_review':syllabus,'institutional_review':inst,'qm_review':qm,'udl_review':udl,'five_e_review':five,'objectives_review':objectives,'citation_review':citations,'compliance_review':compliance,'module_review':modules,'action_plan':actions}

@app.get('/api/health')
def health(): return {'ok':True,'engine':'local-rules','external_ai':False}

@app.post('/api/evaluate')
async def evaluate(course_code:str=Form(...),course_name:str=Form(...),credits:str=Form(''),citation_style:str=Form('APA 7'),course_description:str=Form(''),frameworks:str=Form('[]'),files:List[UploadFile]=File(...)):
    docs=[]
    for f in files:
        raw=await f.read()
        if len(raw)>35*1024*1024: raise HTTPException(400,f'{f.filename} excede 35 MB.')
        try: text=extract(f.filename,raw)
        except Exception as e: text=f'[ERROR DE EXTRACCIÓN: {e}]'
        docs.append({'name':f.filename,'text':text})
    meta={'course_code':course_code,'course_name':course_name,'credits':credits,'citation_style':citation_style,'course_description':course_description,'frameworks':json.loads(frameworks or '[]')}
    data=evaluate_rules(meta,docs); eid=str(uuid.uuid4()); data['evaluation_id']=eid; data['course']={'code':course_code,'name':course_name,'citation_style':citation_style}
    (REPORTS/f'{eid}.json').write_text(json.dumps(data,ensure_ascii=False,indent=2),encoding='utf-8'); return data

@app.get('/api/report/{eid}/docx')
def docx_report(eid:str):
    p=REPORTS/f'{eid}.json'
    if not p.exists(): raise HTTPException(404,'Informe no encontrado.')
    d=json.loads(p.read_text(encoding='utf-8')); out=REPORTS/f'{eid}.docx'; doc=Document(); c=d.get('course',{}); s=d.get('summary',{})
    doc.add_heading('UCAN Course Evaluator',0); doc.add_paragraph(f'{c.get("code","")} – {c.get("name","")}'); doc.add_heading('Resumen ejecutivo',1); doc.add_paragraph(s.get('executive_summary','')); doc.add_paragraph(f'Puntuación: {s.get("overall_score","")}/100'); doc.add_paragraph('Recomendación: '+s.get('recommendation',''))
    sections=[('Prontuario y alineación','syllabus_review'),('Normas institucionales','institutional_review'),('Quality Matters — alto nivel','qm_review'),('DUA / UDL','udl_review'),('Modelo 5E','five_e_review'),('Bloom + Webb DOK','objectives_review'),('Citas y referencias','citation_review'),('Accesibilidad y derechos de autor','compliance_review'),('Módulos','module_review')]
    for title,key in sections:
        doc.add_heading(title,1)
        for x in d.get(key,[]):
            doc.add_heading(x.get('criterion','Hallazgo'),2); doc.add_paragraph(f'{x.get("status","")} | Prioridad: {x.get("priority","")} | Ubicación: {x.get("location","")}')
            if x.get('evidence'): doc.add_paragraph('Evidencia: '+x['evidence'])
            if x.get('comment'): doc.add_paragraph(x['comment'])
            if x.get('recommendation'): doc.add_paragraph('Recomendación: '+x['recommendation'])
            if x.get('example'): doc.add_paragraph('Ejemplo: '+x['example'])
    doc.add_heading('Plan de mejoramiento',1)
    for x in d.get('action_plan',[]): doc.add_paragraph(f'{x.get("priority","")} — {x.get("area","")}: {x.get("action","")} | {x.get("location","")}')
    doc.save(out); return FileResponse(out,filename='Informe_'+re.sub(r'[^A-Za-z0-9_-]+','_',c.get('code','curso'))+'.docx')

app.mount('/',StaticFiles(directory=BASE,html=True),name='static')
